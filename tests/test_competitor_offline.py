"""Offline unit tests for the competitor-analysis agent (no AWS calls).

Run:  ./.venv/bin/python tests/test_competitor_offline.py
"""

from __future__ import annotations

import json
import os
import sys
import tempfile
import time

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from competitor_analysis import corpus
from competitor_analysis.handler import _is_stalled, _parse
from competitor_analysis.models import (
    DIMENSIONS,
    CompetitorAnalysis,
    DimensionFinding,
    EvidenceItem,
    ProcurementDigest,
)
from competitor_analysis.pipeline import _findings_from_synthesis
from competitor_analysis.render_docx import render_docx
from competitor_analysis.render_json import render_json
from competitor_analysis.render_pptx import render_pptx


def test_parse():
    verb, args = _parse('analyze competitor="HP (& EDS)" procurement="CA X (2019)" focal=Deloitte')
    assert verb == "analyze" and args == {
        "competitor": "HP (& EDS)", "procurement": "CA X (2019)", "focal": "Deloitte"}

    verb, args = _parse("analyze competitor=Accenture")
    assert verb == "analyze" and args["competitor"] == "Accenture"
    assert args["procurement"] is None and args["focal"] == "IBM"

    verb, args = _parse("analyze Accenture")   # bare form
    assert verb == "analyze" and args["competitor"] == "Accenture"

    assert _parse("status abc123") == ("status", {"job_id": "abc123"})
    assert _parse("result abc123") == ("result", {"job_id": "abc123"})
    assert _parse("competitors")[0] == "competitors"
    assert _parse("list vendors")[0] == "competitors"
    assert _parse("")[0] == "help"
    assert _parse("analyze")[0] == "help"
    assert _parse("what is this")[0] == "help"
    print("[PASS] parse")


def test_stalled():
    now = time.time()
    assert not _is_stalled({"status": "done"})
    assert not _is_stalled({"status": "running", "worker_deadline_epoch": now + 500})
    assert _is_stalled({"status": "running", "worker_deadline_epoch": now - 500})
    old = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime(now - 4000))
    assert _is_stalled({"status": "running", "updated_at": old})
    print("[PASS] stalled detection")


def test_priority_and_normalize():
    keys = [
        "Acme/P1/random appendix.pdf",
        "Acme/P1/Pricing Workbook.pdf",
        "Acme/P1/Technical Proposal.pdf",
        "Acme/P1/Key Personnel Resumes.pdf",
    ]
    ranked = sorted(keys, key=corpus._priority)
    assert ranked[0].endswith("Pricing Workbook.pdf")
    assert ranked[1].endswith("Technical Proposal.pdf")   # "proposal" outranks "technical"
    assert ranked[-1].endswith("random appendix.pdf")

    assert corpus._normalize("HP (& EDS)") == "hp eds"
    assert corpus._normalize("  Accenture ") == "accenture"
    print("[PASS] priority + normalize")


def _sample_analysis() -> CompetitorAnalysis:
    return CompetitorAnalysis(
        competitor="Acme Corp",
        focal="IBM",
        executive_summary="Acme bids aggressively on price and reuses a COTS platform.",
        dimensions=[
            DimensionFinding(
                key=key, title=spec["title"],
                analysis=f"Narrative for {key}.\n\nSecond paragraph.",
                evidence=[EvidenceItem(procurement="State X ERP", detail="Cited fact.")],
                ibm_implications="Counter by emphasizing delivery record.",
            )
            for key, spec in DIMENSIONS.items()
        ],
        procurement_digests=[
            ProcurementDigest(procurement="State X ERP", client="State X", year="2020",
                              outcome="won", source_docs=["prop.pdf"])
        ],
        docs_analyzed=3,
        warnings=["one procurement skipped"],
    )


def test_renderers():
    analysis = _sample_analysis()
    with tempfile.TemporaryDirectory() as td:
        jpath = render_json(analysis, os.path.join(td, "a.json"))
        data = json.load(open(jpath))
        assert data["competitor"] == "Acme Corp"
        assert len(data["dimensions"]) == len(DIMENSIONS)
        assert data["dimensions"][0]["evidence"][0]["procurement"] == "State X ERP"

        dpath = render_docx(analysis, os.path.join(td, "a.docx"))
        from docx import Document

        texts = "\n".join(p.text for p in Document(dpath).paragraphs)
        assert "Acme Corp" in texts and "Implications for IBM" in texts
        assert "Notes & Limitations" in texts

        ppath = render_pptx(analysis, os.path.join(td, "a.pptx"))
        from pptx import Presentation

        prs = Presentation(ppath)
        deck_text = "\n".join(
            shp.text_frame.text
            for slide in prs.slides
            for shp in slide.shapes
            if shp.has_text_frame
        )
        assert "Acme Corp" in deck_text
        assert "Competitor Bid-Strategy Analysis" in deck_text
        assert "Executive Summary" in deck_text
        assert "Implications for IBM" in deck_text
        # title + summary + glance + 5 dimensions + appendix + notes
        assert len(prs.slides._sldIdLst) >= 9
    print("[PASS] renderers")


def _flowed_boxes_by_slide(prs):
    """Content boxes (below the title band) per slide, as (top, bottom) inches."""
    out = []
    for slide in prs.slides:
        boxes = []
        for shp in slide.shapes:
            top = shp.top / 914400.0
            height = (shp.height or 0) / 914400.0
            if top >= 1.4:                      # skip accent bar / title / subtitle
                boxes.append((top, top + height))
        out.append(sorted(boxes))
    return out


def test_pptx_no_overlap():
    """Long content must paginate; flowed boxes must never overlap or run off-slide."""
    from pptx import Presentation

    long_para = ("This vendor consistently leads with a platform-first message. " * 40)
    analysis = CompetitorAnalysis(
        competitor="Verbose Vendor LLC",
        focal="IBM",
        executive_summary="\n\n".join([long_para, long_para]),
        dimensions=[
            DimensionFinding(
                key=key, title=spec["title"],
                analysis=long_para + "\n\n" + long_para,
                evidence=[
                    EvidenceItem(procurement=f"Procurement {n}",
                                 detail="A specific, sourced fact that is fairly long. " * 4)
                    for n in range(12)
                ],
                ibm_implications=("IBM should counter aggressively across every front. " * 12),
            )
            for key, spec in DIMENSIONS.items()
        ],
        procurement_digests=[
            ProcurementDigest(procurement=f"Procurement {n}", client=f"State {n}",
                              year="2021", outcome="won", source_docs=["a.pdf", "b.pdf"])
            for n in range(30)
        ],
        docs_analyzed=60,
        warnings=[f"warning number {n} about coverage that is reasonably long" for n in range(8)],
    )
    with tempfile.TemporaryDirectory() as td:
        ppath = render_pptx(analysis, os.path.join(td, "long.pptx"))
        prs = Presentation(ppath)

        # Pagination kicked in: far more slides than the base structure.
        assert len(prs.slides._sldIdLst) > 12

        for si, boxes in enumerate(_flowed_boxes_by_slide(prs)):
            prev_bottom = 0.0
            for top, bottom in boxes:
                # no flowed box may run off the bottom of the 7.5" slide
                assert bottom <= 7.5 + 1e-6, f"slide {si}: box bottom {bottom} off-slide"
                # boxes are stacked: each starts at/after the previous one's bottom
                assert top >= prev_bottom - 0.03, (
                    f"slide {si}: overlap (top {top} < prev bottom {prev_bottom})")
                prev_bottom = bottom
    print("[PASS] pptx pagination + no overlap")


def test_findings_from_synthesis():
    raw = {"dimensions": {
        "pricing": {"analysis": "Cheap.", "evidence": [
            {"procurement": "P1", "detail": "Bid $1"}], "ibm_implications": "Undercut."},
    }}
    findings = _findings_from_synthesis(raw)
    assert [f.key for f in findings] == list(DIMENSIONS)
    by_key = {f.key: f for f in findings}
    assert by_key["pricing"].evidence[0].detail == "Bid $1"
    # missing dimensions fail soft with a placeholder, never KeyError
    assert "No evidence" in by_key["staffing"].analysis

    # a truthy NON-dict dimension value (model returned a bare string / list /
    # number) must fail soft too, not crash the whole job
    raw2 = {"dimensions": {
        "solutioning": "just a string, not an object",
        "pricing": ["a", "list"],
        "staffing": 42,
    }}
    findings2 = _findings_from_synthesis(raw2)
    assert [f.key for f in findings2] == list(DIMENSIONS)
    assert all("No evidence" in f.analysis for f in findings2)
    print("[PASS] synthesis mapping")


if __name__ == "__main__":
    test_parse()
    test_stalled()
    test_priority_and_normalize()
    test_renderers()
    test_pptx_no_overlap()
    test_findings_from_synthesis()
    print("ALL PASS")
