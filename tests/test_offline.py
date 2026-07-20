"""Offline component tests — no AWS/Bedrock required.

Exercises models -> rubric -> aggregate -> renderers on a synthetic scorecard
plus a gate-disqualification case, and fills the real PPTX template.

Run directly:  ./.venv/bin/python tests/test_offline.py
(or with pytest if installed).
"""

from __future__ import annotations

import json
import os
import sys
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from scoring_agent.aggregate import aggregate
from scoring_agent.deck_content import (
    CategoryComparison,
    CategoryRow,
    DeckContent,
    DriverRow,
    OutcomeDrivers,
    ProcurementMeta,
    build_category_comparison,
    build_outcome_drivers,
    build_procurement_meta,
)
from scoring_agent.ingest import Document, DocumentText
from scoring_agent.models import (
    Dimension,
    Evidence,
    Provenance,
    ScaleSpec,
    ScoreCell,
    ScorecardResult,
    ScoreMethod,
    SchemeSpec,
    VendorResult,
)
from scoring_agent.render_excel import render_excel
from scoring_agent.render_json import render_json
from scoring_agent.render_pptx import render_pptx
from scoring_agent.rubric import load_base_rubric

OUT = tempfile.mkdtemp(prefix="scoring_test_")


def _cell(vendor, dim_id, value, prov=Provenance.GENERATED, rationale="", ev=None):
    return ScoreCell(
        vendor=vendor,
        dimension_id=dim_id,
        value=value,
        provenance=prov,
        confidence=0.8 if prov == Provenance.GENERATED else None,
        rationale=rationale,
        evidence=ev or [],
    )


def _synthetic_result() -> ScorecardResult:
    scheme = load_base_rubric("demo-project")
    dims = [d.id for d in scheme.dimensions]

    # scripted scores: Deloitte strong, IBM mid, RedMane weaker
    profile = {
        "Deloitte": 5,
        "IBM": 4,
        "RedMane Technology": 3,
    }
    vendors = []
    for name, base in profile.items():
        cells = {}
        for i, did in enumerate(dims):
            # vary a bit so there are strengths/weaknesses
            val = max(1, min(5, base - (1 if i % 3 == 0 else 0)))
            prov = Provenance.EXTRACTED if i % 4 == 0 else Provenance.GENERATED
            ev = (
                [Evidence(doc=f"{name}_proposal.pdf", page=10 + i, quote="…relevant excerpt…")]
                if prov == Provenance.GENERATED
                else []
            )
            cells[did] = _cell(name, did, val, prov, rationale=f"{name} rationale for {did}.", ev=ev)
        vendors.append(VendorResult(vendor=name, cells=cells))

    # give IBM a clear edge on one dim (strength) and a clear deficit on another
    vendors[1].cells["cost_price"].value = 5
    vendors[1].cells["staffing_key_personnel"].value = 2
    return ScorecardResult(
        project_id="demo-project", focal_vendor="IBM", scheme=scheme, vendors=vendors
    )


def test_aggregate_and_ci():
    result = aggregate(_synthetic_result())
    for v in result.vendors:
        assert v.normalized_total_pct is not None, f"{v.vendor} has no total"
        for cell in v.cells.values():
            assert cell.normalized_pct is not None
    # Deloitte (all 5s minus dips) should rank #1
    top = result.ranked_vendors()[0]
    assert top.vendor == "Deloitte", f"expected Deloitte top, got {top.vendor}"
    # CI centered on IBM
    assert result.ci is not None
    assert result.ci.focal_vendor == "IBM"
    ibm = next(v for v in result.vendors if v.vendor == "IBM")
    assert result.ci.rank == ibm.rank
    assert result.ci.summary
    assert "staffing" in " ".join(result.ci.weaknesses).lower() or result.ci.key_drivers
    print("  aggregate+CI:", result.ci.summary)
    return result


def test_gate_disqualification():
    scheme = SchemeSpec(
        project_id="gate-test",
        method=ScoreMethod.PASS_FAIL_PLUS_POINTS,
        scale=ScaleSpec(min=1, max=5),
        dimensions=[
            Dimension(id="min_req", name="Minimum Requirements", mandatory_gate=True, weight=0.0),
            Dimension(id="tech", name="Technical", weight=0.6),
            Dimension(id="cost", name="Cost", weight=0.4, cost_related=True),
        ],
    )
    good = VendorResult(
        vendor="GoodCo",
        cells={
            "min_req": _cell("GoodCo", "min_req", 5, Provenance.GENERATED),
            "tech": _cell("GoodCo", "tech", 4, Provenance.GENERATED),
            "cost": _cell("GoodCo", "cost", 3, Provenance.GENERATED),
        },
    )
    fail = VendorResult(
        vendor="FailCo",
        cells={
            "min_req": ScoreCell("FailCo", "min_req", provenance=Provenance.GATE_FAIL,
                                 rationale="missing required certification"),
            "tech": _cell("FailCo", "tech", 5, Provenance.GENERATED),
            "cost": _cell("FailCo", "cost", 5, Provenance.GENERATED),
        },
    )
    result = aggregate(
        ScorecardResult(project_id="gate-test", focal_vendor="GoodCo", scheme=scheme,
                        vendors=[good, fail])
    )
    fc = next(v for v in result.vendors if v.vendor == "FailCo")
    gc = next(v for v in result.vendors if v.vendor == "GoodCo")
    assert fc.disqualified and fc.rank is None, "FailCo should be disqualified"
    assert fc.disqualification_reason
    assert gc.rank == 1, "GoodCo should rank #1"
    print("  gate DQ reason:", fc.disqualification_reason)


def test_json_roundtrip():
    result = aggregate(_synthetic_result())
    path = os.path.join(OUT, "scorecard.json")
    text = render_json(result, path)
    reloaded = ScorecardResult.from_dict(json.loads(text))
    assert reloaded.project_id == result.project_id
    assert reloaded.focal_vendor == result.focal_vendor
    assert len(reloaded.vendors) == len(result.vendors)
    assert reloaded.ci is not None and reloaded.ci.focal_vendor == "IBM"
    # provenance + evidence survive the round-trip
    v = reloaded.vendors[0]
    any_ev = any(c.evidence for c in v.cells.values())
    assert any_ev, "evidence lost in round-trip"
    print("  json bytes:", len(text), "->", path)


def test_render_excel():
    result = aggregate(_synthetic_result())
    path = os.path.join(OUT, "scorecard.xlsx")
    render_excel(result, path)
    assert os.path.getsize(path) > 0
    from openpyxl import load_workbook

    wb = load_workbook(path)
    ws = wb["Scorecard"]
    assert ws.cell(row=1, column=1).value == "Evaluation Dimension"
    headers = [ws.cell(row=1, column=c).value for c in range(3, 3 + len(result.vendors))]
    assert "IBM" in headers, headers
    assert "Evidence & Notes" in wb.sheetnames and "Scheme" in wb.sheetnames
    print("  excel:", path, "| vendor headers:", headers)


def test_technical_financial_split():
    r = _synthetic_result()
    # unambiguous cost ordering: IBM strongest, Deloitte weakest on price
    for v in r.vendors:
        if v.vendor == "IBM":
            v.cells["cost_price"].value = 5
        elif v.vendor == "Deloitte":
            v.cells["cost_price"].value = 2
    result = aggregate(r)
    # base rubric has a cost_related dim (cost_price) -> both splits populated
    for v in result.vendors:
        assert v.technical_pct is not None, f"{v.vendor} missing technical_pct"
        assert v.financial_pct is not None, f"{v.vendor} missing financial_pct"
    ibm = next(v for v in result.vendors if v.vendor == "IBM")
    dl = next(v for v in result.vendors if v.vendor == "Deloitte")
    # the split differentiates: Deloitte leads technical, IBM leads financial
    assert dl.technical_rank == 1, dl.technical_rank
    assert ibm.financial_rank == 1, ibm.financial_rank
    assert dl.financial_rank != 1, dl.financial_rank
    print("  split: financial#1=IBM technical#1=Deloitte")


def _synthetic_deck_content() -> DeckContent:
    return DeckContent(
        meta=ProcurementMeta(
            agency="State Demo Department of Services",
            rfp_number="RFQ-DEMO-2025-01",
            tcv="$5,800,000",
            summary="Implement and operate a demo case-management system.",
            winning_vendor="Deloitte",
            vendors=["Deloitte", "IBM", "RedMane Technology"],
            documents=["rfp.pdf", "deloitte_proposal.pdf", "scoresheet.pdf"],
        ),
        drivers=OutcomeDrivers(
            winner="Deloitte", focal="IBM",
            why_won=[DriverRow("Narrative excellence", "480/550 narrative", "Deep understanding")],
            why_focal_lost=[DriverRow("Narrative weakness", "237.5 vs 480", "Largest scoring gap")],
        ),
        comparison=CategoryComparison(
            focal="IBM", winner="Deloitte",
            rows=[
                CategoryRow("Price", ["$5.8M fixed"], ["Mid-range, competitive"]),
                CategoryRow("Testing", ["SIT + hypercare"], ["Parallel payroll testing"]),
            ],
        ),
    )


def test_render_deck():
    result = aggregate(_synthetic_result())
    content = _synthetic_deck_content()
    path = os.path.join(OUT, "scorecard.pptx")
    render_pptx(result, path, deck_content=content)
    assert os.path.getsize(path) > 0
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(path)
    slide_w, slide_h = prs.slide_width, prs.slide_height
    # The deck paginates large tables across continuation slides, so the count
    # grows with dimensions/vendors — but must always cover the core sections.
    assert len(prs.slides) >= 7, len(prs.slides)

    def slide_text(slide):
        parts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                parts.append(sh.text_frame.text)
            if getattr(sh, "has_table", False):
                for row in sh.table.rows:
                    for c in row.cells:
                        parts.append(c.text)
        return " ".join(parts)

    all_text = "\n".join(slide_text(s) for s in prs.slides)
    for expected in ("Overview", "Final Scoring", "Scoring Overview", "Detailed Scoring",
                     "Outcome Drivers", "Why Deloitte Won", "IBM vs. Deloitte"):
        assert expected in all_text, f"missing slide/section: {expected}"
    # branding + metadata + narrative made it in
    assert "SLED Competitive Intelligence" in all_text, "footer/branding lost"
    assert "RFQ-DEMO-2025-01" in all_text and "$5,800,000" in all_text
    assert "Very strong" in all_text or "Weak" in all_text, "RAG adjectives missing"

    # COMPLETENESS: every dimension, every vendor, and the totals rows must
    # actually appear somewhere in the deck (the bug this guards against was
    # rows being clipped off the bottom of an over-stuffed slide).
    for dim in result.scheme.dimensions:
        assert dim.name in all_text, f"dimension clipped from deck: {dim.name}"
    for v in result.vendors:
        assert v.vendor in all_text, f"vendor clipped from deck: {v.vendor}"
    for totals in ("Normalized Total", "Rank"):
        assert totals in all_text, f"totals row missing from deck: {totals}"

    # NO OVERFLOW: no table may extend past the visible slide area, or its lower
    # rows are silently clipped when the deck is opened/presented.
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if getattr(sh, "has_table", False):
                assert sh.left + sh.width <= slide_w + Inches(0.05), f"slide {si}: table over right edge"
                assert sh.top + sh.height <= slide_h + Inches(0.05), f"slide {si}: table over bottom edge"
    print("  deck:", path, "| slides:", len(prs.slides))


class _MockDeckBedrock:
    """Canned Bedrock replies for the deck-content builders."""

    def converse(self, user_text, **kwargs):
        return ""

    def converse_json(self, user_text, **kwargs):
        if "SOLICITATION TEXT" in user_text:
            return {"agency": "Demo Agency", "rfp_number": "RFQ-9",
                    "tcv": "$1,000,000", "summary": "A demo procurement."}
        if "why_won" in user_text:
            return {"why_won": [{"factor": "Better narrative", "evidence": "higher score",
                                 "impact": "won phase 1"}],
                    "why_focal_lost": [{"factor": "Narrative gap", "evidence": "trailed leader",
                                        "impact": "lower rank"}]}
        if "PROPOSAL EXCERPTS" in user_text:
            return {"rows": [{"category": "Price", "focal_points": ["fixed price"],
                              "winner_points": ["lower cost"]}]}
        raise AssertionError("unexpected deck prompt:\n" + user_text[:150])


def test_deck_content_mock():
    result = aggregate(_synthetic_result())
    mock = _MockDeckBedrock()
    warnings = []

    meta = build_procurement_meta(result, None, "Solicitation text here.", bedrock=mock, warnings=warnings)
    assert meta.rfp_number == "RFQ-9" and meta.agency == "Demo Agency"
    assert meta.winning_vendor == "Deloitte"  # derived from ranking (top of synthetic)

    drivers = build_outcome_drivers(result, bedrock=mock, warnings=warnings)
    assert drivers is not None and drivers.winner == "Deloitte" and drivers.focal == "IBM"
    assert drivers.why_won and drivers.why_focal_lost

    proposal_texts = {
        "IBM": [DocumentText(Document(category="proposal", filename="ibm.pdf", vendor="IBM"),
                             pages=["IBM proposal narrative about price and testing."])],
        "Deloitte": [DocumentText(Document(category="proposal", filename="dl.pdf", vendor="Deloitte"),
                                  pages=["Deloitte proposal narrative about price and testing."])],
    }
    comp = build_category_comparison(result, proposal_texts, bedrock=mock, warnings=warnings)
    assert comp is not None and comp.rows and comp.rows[0].category == "Price"
    assert not warnings, warnings
    print("  deck content: meta+drivers+comparison OK")


def main():
    tests = [
        test_aggregate_and_ci,
        test_gate_disqualification,
        test_technical_financial_split,
        test_json_roundtrip,
        test_render_excel,
        test_render_deck,
        test_deck_content_mock,
    ]
    failures = 0
    for t in tests:
        try:
            print(f"[RUN] {t.__name__}")
            t()
            print(f"[PASS] {t.__name__}")
        except Exception as e:  # noqa: BLE001
            failures += 1
            import traceback

            print(f"[FAIL] {t.__name__}: {e}")
            traceback.print_exc()
    print(f"\nOutputs in: {OUT}")
    print("DONE" if failures == 0 else f"{failures} FAILURE(S)")
    sys.exit(1 if failures else 0)


if __name__ == "__main__":
    main()
