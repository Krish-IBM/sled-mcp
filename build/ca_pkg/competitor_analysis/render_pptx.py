"""PowerPoint deck renderer — IBM-branded competitor bid-strategy deck.

Renders the same ``CompetitorAnalysis`` as the Word report into a clean 16:9
deck (python-pptx, no template file). Structure:

    1. Title
    2. Executive Summary            (auto-paginated if long)
    3. Strategy at a Glance         (5-dimension headline table)
    4..N. One section per dimension — narrative + evidence + implications,
          each block flowed top-to-bottom and spilled onto continuation slides
          so text NEVER overlaps regardless of length.
    N+1. Appendix: Procurements Analyzed
    N+2. Notes & Limitations        (only if warnings exist)

The overlap guarantee comes from a small flow engine (``_Flow``): every block's
height is estimated, blocks are stacked with a fixed gap, a new slide is started
before the cursor would cross the footer, and each text box also carries
``TEXT_TO_FIT_SHAPE`` autofit as a belt-and-suspenders so any estimate shortfall
shrinks the text instead of overflowing into the next block.
"""

from __future__ import annotations

import math
from typing import List, Optional

from .models import CompetitorAnalysis, DimensionFinding

# ── palette / branding (matches scoring_agent deck) ──────────────────────────
_FONT = "IBM Plex Sans"
_IBM_BLUE = (0x0F, 0x62, 0xFE)
_IBM_DARK = (0x00, 0x2D, 0x9C)
_WHITE = (0xFF, 0xFF, 0xFF)
_INK = (0x16, 0x16, 0x16)
_GRAY = (0x6F, 0x6F, 0x6F)
_LABEL_BG = (0xF2, 0xF4, 0xF8)
_ROW_ALT = (0xF7, 0xF8, 0xFA)
_HILITE = (0xE8, 0xF0, 0xFE)   # implications callout fill

# ── 16:9 canvas geometry (inches) ────────────────────────────────────────────
_SLIDE_W = 13.333
_SLIDE_H = 7.5
_MARGIN = 0.6
_CONTENT_W = _SLIDE_W - 2 * _MARGIN
_BODY_TOP = 1.45
_BODY_BOTTOM = 6.95            # keep clear of the footer at 7.15
_GAP = 0.10
_FOOTER = "Prepared by IBM SLED Competitive Intelligence"

# ── text-height estimation ───────────────────────────────────────────────────
# Conservative on purpose: over-estimating height only opens whitespace, while
# under-estimating is what causes overlap. ~130 chars per (inch * pt / size).
_CPI = 130.0
_LINE_H = 1.32                 # line-height multiple


def _rgb(t):
    from pptx.dml.color import RGBColor

    return RGBColor(*t)


def _est_lines(text: str, width_in: float, size: float, hanging: float = 0.0) -> int:
    usable = max(0.5, width_in - hanging)
    cpl = max(4, int(usable * _CPI / size))
    lines = 0
    for seg in (text or " ").split("\n"):
        lines += max(1, math.ceil(len(seg or " ") / cpl))
    return lines


def _text_height(text: str, width_in: float, size: float, *, hanging: float = 0.0,
                 pad: float = 0.14) -> float:
    lines = _est_lines(text, width_in, size, hanging)
    return lines * size * _LINE_H / 72.0 + pad


# ── low-level slide helpers ───────────────────────────────────────────────────
def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # layout 6 = "Blank"


def _autofit(tf):
    """Shrink text to fit its box if an estimate ever falls short (no overflow)."""
    from pptx.enum.text import MSO_AUTO_SIZE

    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:  # noqa: BLE001 — never let autofit config break rendering
        pass


def _textbox(slide, text, left, top, width, height, *, size=14, bold=False,
             color=_INK, align=None, anchor=None, autofit=False):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor if anchor is not None else MSO_ANCHOR.TOP
    if autofit:
        _autofit(tf)
    for i, line in enumerate((text or "").split("\n") or [""]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align if align is not None else PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line
        run.font.name = _FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
    return box


def _accent_bar(slide):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.24), Inches(_SLIDE_H))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(_IBM_BLUE)
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _footer(slide):
    _textbox(slide, _FOOTER, _MARGIN, _SLIDE_H - 0.35, _CONTENT_W, 0.3, size=9, color=_GRAY)


def _title_block(slide, title: str, subtitle: str = ""):
    _accent_bar(slide)
    _textbox(slide, title, _MARGIN, 0.42, _CONTENT_W, 0.75, size=26, bold=True, color=_IBM_DARK)
    if subtitle:
        _textbox(slide, subtitle, _MARGIN, 1.08, _CONTENT_W, 0.32, size=12, color=_GRAY)
    _footer(slide)


# ── content-flow engine (guarantees no overlap) ───────────────────────────────
class _Flow:
    """Stacks blocks top-to-bottom, spilling onto continuation slides.

    Each block is placed at the running cursor and the cursor advances by the
    block's height + gap. Before placing, if the block will not fit and we are
    not already at the top of a fresh slide, a continuation slide is started.
    A block taller than a whole page is clamped to the page and relies on
    text-autofit — this only happens for pathologically long single paragraphs.
    """

    def __init__(self, prs, title: str, subtitle: str = ""):
        self.prs = prs
        self.title = title
        self.subtitle = subtitle
        self._start(first=True)

    def _start(self, first: bool):
        self.slide = _blank_slide(self.prs)
        _title_block(self.slide, self.title if first else f"{self.title}  (cont.)",
                     self.subtitle if first else "")
        self.y = _BODY_TOP

    def _fit(self, height: float, keep_with: float = 0.0):
        """Ensure `height` (+ optional following block) fits; else new slide."""
        need = height + keep_with
        if self.y + need > _BODY_BOTTOM and self.y > _BODY_TOP + 1e-6:
            self._start(first=False)

    def _avail(self) -> float:
        return _BODY_BOTTOM - self.y

    def heading(self, text: str, *, size=14, gap_before=0.12):
        h = _text_height(text, _CONTENT_W, size, pad=0.06)
        # Keep a heading with at least the first line of what follows it.
        self._fit(h + gap_before, keep_with=0.55)
        self.y += gap_before
        _textbox(self.slide, text, _MARGIN, self.y, _CONTENT_W, h,
                 size=size, bold=True, color=_IBM_DARK)
        self.y += h + 0.04

    def paragraph(self, text: str, *, size=12, color=_INK):
        if not (text or "").strip():
            return
        h = _text_height(text, _CONTENT_W, size)
        self._fit(h)
        placed = min(h, self._avail())
        _textbox(self.slide, text, _MARGIN, self.y, _CONTENT_W, placed,
                 size=size, color=color, autofit=True)
        self.y += placed + _GAP

    def bullet(self, lead: str, text: str, *, size=11):
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        combined = (f"{lead}: {text}" if lead else text) or ""
        h = _text_height("• " + combined, _CONTENT_W, size, hanging=0.28, pad=0.10)
        self._fit(h)
        placed = min(h, self._avail())
        box = self.slide.shapes.add_textbox(
            Inches(_MARGIN), Inches(self.y), Inches(_CONTENT_W), Inches(placed))
        tf = box.text_frame
        tf.word_wrap = True
        _autofit(tf)
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT

        def _run(t, *, bold=False, color=_INK):
            r = para.add_run()
            r.text = t
            r.font.name = _FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = _rgb(color)
            return r

        _run("•  ", color=_IBM_BLUE, bold=True)
        if lead:
            _run(f"{lead}: ", bold=True, color=_IBM_DARK)
        _run(text)
        self.y += placed + 0.05

    def callout(self, heading: str, body: str, *, size=11):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR
        from pptx.util import Inches, Pt

        if not (body or "").strip():
            return
        inner_w = _CONTENT_W - 0.4
        h = (_text_height(heading, inner_w, size + 1, pad=0.02)
             + _text_height(body, inner_w, size, pad=0.02) + 0.34)
        self._fit(h, keep_with=0.0)
        placed = min(h, self._avail())
        shp = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(_MARGIN), Inches(self.y),
            Inches(_CONTENT_W), Inches(placed))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(_HILITE)
        shp.line.color.rgb = _rgb(_IBM_BLUE)
        shp.line.width = Pt(1)
        shp.shadow.inherit = False
        tf = shp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.18)
        tf.margin_top = Inches(0.12)
        tf.margin_bottom = Inches(0.12)
        _autofit(tf)

        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = heading
        r0.font.name = _FONT
        r0.font.size = Pt(size + 1)
        r0.font.bold = True
        r0.font.color.rgb = _rgb(_IBM_DARK)

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = body
        r1.font.name = _FONT
        r1.font.size = Pt(size)
        r1.font.color.rgb = _rgb(_INK)
        self.y += placed + _GAP


# ── table helper (glance slide) ───────────────────────────────────────────────
def _set_cell(cell, text, *, fill=None, color=_INK, bold=False, size=11, align=None):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    cell.vertical_anchor = MSO_ANCHOR.TOP
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(fill)
    tf = cell.text_frame
    tf.word_wrap = True
    for i, line in enumerate((text or "").split("\n") or [""]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align if align is not None else PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line
        run.font.name = _FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)


def _first_sentence(text: str, limit: int = 200) -> str:
    t = " ".join((text or "").split())
    if not t:
        return "—"
    dot = t.find(". ")
    if 0 < dot < limit:
        return t[: dot + 1]
    return t if len(t) <= limit else t[: limit - 1].rstrip() + "…"


# ── slides ────────────────────────────────────────────────────────────────────
def _slide_title(prs, a: CompetitorAnalysis):
    slide = _blank_slide(prs)
    _accent_bar(slide)
    _textbox(slide, a.competitor, _MARGIN, 2.35, _CONTENT_W, 1.3, size=40, bold=True,
             color=_IBM_DARK)
    _textbox(slide, "Competitor Bid-Strategy Analysis", _MARGIN, 3.75, _CONTENT_W, 0.6,
             size=20, color=_INK)
    meta = (f"Focal vendor: {a.focal}      |      "
            f"Procurements analyzed: {len(a.procurement_digests)}      |      "
            f"Documents analyzed: {a.docs_analyzed}")
    _textbox(slide, meta, _MARGIN, 4.5, _CONTENT_W, 0.4, size=13, color=_GRAY)
    _textbox(slide, f"Generated {a.generated_at}", _MARGIN, 4.9, _CONTENT_W, 0.35,
             size=11, color=_GRAY)
    _textbox(slide,
             "Source: FOIA documents (proposals, pricing workbooks, evaluation scoresheets) "
             "from the SLED competitive-intelligence corpus.",
             _MARGIN, 5.35, _CONTENT_W, 0.6, size=11, color=_GRAY)
    _footer(slide)


def _slide_executive_summary(prs, a: CompetitorAnalysis):
    flow = _Flow(prs, "Executive Summary",
                 f"{a.competitor}  ·  focal vendor {a.focal}")
    summary = a.executive_summary or "No summary produced."
    # Split into sentences-ish paragraphs so long summaries paginate cleanly.
    for para in summary.split("\n\n"):
        flow.paragraph(para.strip(), size=14)


def _slide_glance(prs, a: CompetitorAnalysis):
    from pptx.util import Inches

    slide = _blank_slide(prs)
    _title_block(slide, "Strategy at a Glance",
                 "One-line headline per dimension — detail follows")
    dims = a.dimensions or []
    rows = 1 + max(len(dims), 1)
    gf = slide.shapes.add_table(rows, 2, Inches(_MARGIN), Inches(_BODY_TOP),
                                Inches(_CONTENT_W), Inches(min(0.9 * rows, 5.2)))
    table = gf.table
    table.first_row = False
    table.horz_banding = False
    table.columns[0].width = Inches(3.6)
    table.columns[1].width = Inches(_CONTENT_W - 3.6)
    _set_cell(table.cell(0, 0), "Dimension", fill=_IBM_BLUE, color=_WHITE, bold=True, size=13)
    _set_cell(table.cell(0, 1), "Headline", fill=_IBM_BLUE, color=_WHITE, bold=True, size=13)
    if not dims:
        _set_cell(table.cell(1, 0), "—", fill=_LABEL_BG, size=11)
        _set_cell(table.cell(1, 1), "No dimensions produced.", fill=_WHITE, size=11)
        return
    for i, dim in enumerate(dims, start=1):
        fill = _ROW_ALT if i % 2 else _WHITE
        _set_cell(table.cell(i, 0), dim.title, fill=_LABEL_BG, bold=True, size=11)
        _set_cell(table.cell(i, 1), _first_sentence(dim.analysis), fill=fill, size=11)


def _slide_dimension(prs, a: CompetitorAnalysis, index: int, dim: DimensionFinding):
    flow = _Flow(prs, f"{index} · {dim.title}",
                 f"{a.competitor} — bid-strategy analysis")
    for para in (dim.analysis or "No evidence found in the corpus.").split("\n\n"):
        flow.paragraph(para.strip(), size=13)
    if dim.evidence:
        flow.heading("Evidence")
        for item in dim.evidence:
            detail = (item.detail or "").strip()
            if detail:
                flow.bullet(item.procurement.strip(), detail)
    if (dim.ibm_implications or "").strip():
        flow.callout(f"Implications for {a.focal}", dim.ibm_implications.strip())


def _slide_appendix(prs, a: CompetitorAnalysis):
    flow = _Flow(prs, "Appendix: Procurements Analyzed",
                 f"{len(a.procurement_digests)} procurements  ·  "
                 f"{a.docs_analyzed} FOIA documents")
    if not a.procurement_digests:
        flow.paragraph("No procurement digests were produced.", size=12)
        return
    for d in a.procurement_digests:
        extras = ", ".join(x for x in (d.client, d.year, d.outcome) if x)
        tail = []
        if extras:
            tail.append(f"({extras})")
        if d.source_docs:
            tail.append(f"{len(d.source_docs)} docs")
        flow.bullet(d.procurement, " ".join(tail) if tail else "")


def _slide_notes(prs, a: CompetitorAnalysis):
    if not a.warnings:
        return
    flow = _Flow(prs, "Notes & Limitations",
                 "Coverage caveats from this run")
    for w in a.warnings:
        flow.bullet("", w)


# ── entry point ───────────────────────────────────────────────────────────────
def render_pptx(analysis: CompetitorAnalysis, path: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W)
    prs.slide_height = Inches(_SLIDE_H)

    _slide_title(prs, analysis)
    _slide_executive_summary(prs, analysis)
    _slide_glance(prs, analysis)
    for i, dim in enumerate(analysis.dimensions, start=1):
        _slide_dimension(prs, analysis, i, dim)
    _slide_appendix(prs, analysis)
    _slide_notes(prs, analysis)

    prs.save(path)
    return path
