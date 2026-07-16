"""PowerPoint deck renderer — IBM-branded competitive scorecard.

Builds a multi-slide deck programmatically (python-pptx, no template file
required) in the structure of the SLED Competitive Intelligence decks:

    1. Title
    2. Overview            — procurement metadata + documents received
    3. Final Scoring       — Technical Rank / Financial Rank / Final Score
    4. Scoring Overview    — winner summary box + RAG dimension table
    5. Detailed Scoring    — full dimensions x vendors matrix (provenance-colored)
    6. Outcome Drivers     — Why the winner won / Why the focal vendor lost
    7. Category Comparison — focal vs. winner (Price / Implementation / Testing / …)

Slides 2/4/5 render purely from the scored data. Slides 3 uses the aggregate's
technical/financial split. Slides 2 (metadata), 6, and 7 use ``deck_content``
(Bedrock narrative); each degrades gracefully when that content is absent so the
deck always renders. ``template_path`` is accepted for backward compatibility but
the deck is generated on a blank presentation with IBM branding applied here.
"""

from __future__ import annotations

from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .deck_content import DeckContent, resolve_focal, resolve_winner
from .models import Provenance, ScorecardResult, ScoreCell, VendorResult

# ── palette ──────────────────────────────────────────────────────────────── #
_FONT = "IBM Plex Sans"
_IBM_BLUE = RGBColor(0x0F, 0x62, 0xFE)
_IBM_DARK = RGBColor(0x00, 0x2D, 0x9C)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_INK = RGBColor(0x16, 0x16, 0x16)
_LABEL_BG = RGBColor(0xF2, 0xF4, 0xF8)
_ROW_ALT = RGBColor(0xF7, 0xF8, 0xFA)

# provenance fills (detailed matrix)
_FILL = {
    Provenance.EXTRACTED: RGBColor(0xD6, 0xEA, 0xD3),
    Provenance.GENERATED: RGBColor(0xDC, 0xE6, 0xF1),
    Provenance.GATE_FAIL: RGBColor(0xF4, 0xCC, 0xCC),
    Provenance.NOT_SCORED: RGBColor(0xF2, 0xF2, 0xF2),
}
# RAG fills (scoring overview)
_RAG_GREEN = RGBColor(0xC6, 0xE0, 0xB4)
_RAG_AMBER = RGBColor(0xFF, 0xE6, 0x99)
_RAG_RED = RGBColor(0xE6, 0xA0, 0x92)

# 16:9 canvas
_SLIDE_W = 13.333
_SLIDE_H = 7.5
_MARGIN = 0.5
_CONTENT_W = _SLIDE_W - 2 * _MARGIN
_BODY_TOP = 1.35
_FOOTER = "Prepared by SLED Competitive Intelligence Team"


# ── low-level helpers ────────────────────────────────────────────────────── #
def _blank_slide(prs: Presentation):
    # layout 6 is "Blank" in python-pptx's default template
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, text, left, top, width, height, *, size=14, bold=False,
             color=_INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = (text or "").split("\n") or [""]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = _FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def _accent_bar(slide):
    from pptx.enum.shapes import MSO_SHAPE

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.28), Inches(_SLIDE_H))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _IBM_BLUE
    bar.line.fill.background()
    return bar


def _title_block(slide, title: str, subtitle: str = "") -> None:
    _accent_bar(slide)
    _textbox(slide, title, _MARGIN, 0.42, _CONTENT_W, 0.8, size=28, bold=True, color=_IBM_DARK)
    if subtitle:
        _textbox(slide, subtitle, _MARGIN, 1.02, _CONTENT_W, 0.35, size=13, color=_INK)
    _footer(slide)


def _footer(slide) -> None:
    _textbox(slide, _FOOTER, _MARGIN, _SLIDE_H - 0.35, _CONTENT_W, 0.3, size=9, color=RGBColor(0x6F, 0x6F, 0x6F))


def _set_cell(cell, text, *, fill=None, color=_INK, bold=False, size=11, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.MIDDLE):
    cell.vertical_anchor = anchor
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    lines = (text if text else "").split("\n") or [""]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = _FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _table(slide, n_rows, n_cols, left, top, width, height):
    gf = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = gf.table
    table.first_row = False
    table.horz_banding = False
    return table


def _header_row(table, labels, *, size=12):
    for j, label in enumerate(labels):
        align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        _set_cell(table.cell(0, j), label, fill=_IBM_BLUE, color=_WHITE, bold=True, size=size, align=align)


def _bullets(points: List[str]) -> str:
    return "\n".join(f"• {p}" for p in points) if points else "—"


# ── value formatting ─────────────────────────────────────────────────────── #
def _cell_display(cell: Optional[ScoreCell]) -> str:
    if cell is None:
        return "—"
    if cell.provenance == Provenance.GATE_FAIL:
        return "FAIL"
    if cell.native_value:
        return cell.native_value
    if cell.value is not None:
        return f"{cell.value:g}"
    return "—"


def _adjective(pct: Optional[float]) -> str:
    if pct is None:
        return "—"
    if pct >= 80:
        return "Very strong"
    if pct >= 65:
        return "Strong"
    if pct >= 45:
        return "Moderate"
    return "Weak"


def _rag_fill(pct: Optional[float]):
    if pct is None:
        return _LABEL_BG
    if pct >= 65:
        return _RAG_GREEN
    if pct >= 45:
        return _RAG_AMBER
    return _RAG_RED


def _final_score(v: VendorResult) -> str:
    if v.native_total is not None:
        return f"{v.native_total:.0f}"
    if v.normalized_total_pct is not None:
        return f"{v.normalized_total_pct:.0f}%"
    return "—"


def _rank_tag(v: VendorResult) -> str:
    if v.disqualified:
        return "DQ"
    return f"#{v.rank}" if v.rank else "—"


def _clamp(v, lo, hi):
    return max(lo, min(hi, v))


# ── slides ───────────────────────────────────────────────────────────────── #
def _slide_title(prs, result: ScorecardResult, content: DeckContent) -> None:
    slide = _blank_slide(prs)
    _accent_bar(slide)
    meta = content.meta
    heading = meta.agency or result.project_id.replace("_", " ")
    _textbox(slide, heading, _MARGIN, 2.4, _CONTENT_W, 1.6, size=34, bold=True, color=_IBM_DARK)
    _textbox(slide, "Competitive Scorecard & FOIA Analysis", _MARGIN, 3.9, _CONTENT_W, 0.6,
             size=18, color=_INK)
    sub = []
    if meta.rfp_number:
        sub.append(f"Solicitation: {meta.rfp_number}")
    sub.append(f"Focal vendor: {result.focal_vendor}")
    _textbox(slide, "   |   ".join(sub), _MARGIN, 4.6, _CONTENT_W, 0.5, size=13,
             color=RGBColor(0x6F, 0x6F, 0x6F))
    _footer(slide)


def _slide_overview(prs, result: ScorecardResult, content: DeckContent) -> None:
    slide = _blank_slide(prs)
    _title_block(slide, "Overview")
    meta = content.meta
    winner = meta.winning_vendor or (resolve_winner(result).vendor if resolve_winner(result) else "—")
    rows = [
        ("Procurement Summary", meta.summary or "—"),
        ("RFP #", meta.rfp_number or "—"),
        ("Procuring Agency", meta.agency or "—"),
        ("Total Contract Value", meta.tcv or "—"),
        ("Winning Vendor", winner),
        ("List of all Vendors", "\n".join(f"• {v}" for v in (meta.vendors or result.vendor_names())) or "—"),
        ("Documents Received", "\n".join(f"• {d}" for d in meta.documents[:12]) or "—"),
    ]
    table = _table(slide, len(rows), 2, _MARGIN, _BODY_TOP, _CONTENT_W, 5.6)
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(_CONTENT_W - 3.2)
    for r, (label, value) in enumerate(rows):
        _set_cell(table.cell(r, 0), label, fill=_IBM_BLUE, color=_WHITE, bold=True, size=12,
                  anchor=MSO_ANCHOR.TOP)
        _set_cell(table.cell(r, 1), value, fill=(_ROW_ALT if r % 2 else _WHITE), size=11,
                  anchor=MSO_ANCHOR.TOP)


def _slide_final_scoring(prs, result: ScorecardResult, content: DeckContent) -> None:
    slide = _blank_slide(prs)
    _title_block(slide, "Final Scoring")
    vendors = result.ranked_vendors()
    has_financial = any(v.financial_pct is not None for v in vendors)

    cols = ["Vendor", "Technical Rank", "Financial Rank", "Final Score"] if has_financial \
        else ["Vendor", "Technical Rank", "Final Score"]
    focal = resolve_focal(result)
    winner = resolve_winner(result)

    # lead narrative
    if content.drivers and content.drivers.why_won:
        lead = f"{winner.vendor} was selected — {content.drivers.why_won[0].factor}." if winner else ""
    elif result.ci and result.ci.summary:
        lead = result.ci.summary
    else:
        lead = ""
    if lead:
        _textbox(slide, lead, _MARGIN, _BODY_TOP, _CONTENT_W, 0.7, size=12, color=_INK)

    top = _BODY_TOP + (0.8 if lead else 0.05)
    table = _table(slide, 1 + len(vendors), len(cols), _MARGIN, top, _CONTENT_W,
                   min(0.55 * (len(vendors) + 1), 5.4))
    _header_row(table, cols)

    def _ord(n: Optional[int]) -> str:
        return f"{n}" if n else "—"

    for i, v in enumerate(vendors, start=1):
        is_focal = focal is not None and v.vendor == focal.vendor
        base_fill = RGBColor(0xE8, 0xF0, 0xFE) if is_focal else (_ROW_ALT if i % 2 else _WHITE)
        _set_cell(table.cell(i, 0), v.vendor, fill=base_fill, bold=is_focal, size=11)
        _set_cell(table.cell(i, 1), _ord(v.technical_rank), fill=base_fill, size=11, align=PP_ALIGN.CENTER)
        c = 2
        if has_financial:
            _set_cell(table.cell(i, 2), _ord(v.financial_rank), fill=base_fill, size=11, align=PP_ALIGN.CENTER)
            c = 3
        final = "Winner" if (winner and v.vendor == winner.vendor) else _final_score(v)
        _set_cell(table.cell(i, c), final, fill=base_fill, bold=(winner and v.vendor == winner.vendor),
                  size=11, align=PP_ALIGN.CENTER)


def _slide_scoring_overview(prs, result: ScorecardResult, content: DeckContent) -> None:
    slide = _blank_slide(prs)
    _title_block(slide, "Scoring Overview")
    ranked = result.ranked_vendors()
    focal = resolve_focal(result)
    winner = ranked[0] if ranked else None

    # winner summary box (label | value)
    driver = ""
    if content.drivers and content.drivers.why_won:
        driver = "; ".join(r.factor for r in content.drivers.why_won[:2])
    elif result.ci and result.ci.key_drivers:
        driver = "Key gaps: " + ", ".join(result.ci.key_drivers)
    focal_outcome = "—"
    if focal is not None:
        focal_outcome = "Disqualified" if focal.disqualified else (
            f"Rank #{focal.rank} of {len(ranked)}" if focal.rank else "Unranked")
    box_rows = [
        ("Winner", winner.vendor if winner else "—"),
        ("Final Score", _final_score(winner) if winner else "—"),
        ("#2", ranked[1].vendor if len(ranked) > 1 else "—"),
        ("#3", ranked[2].vendor if len(ranked) > 2 else "—"),
        (f"{result.focal_vendor} Outcome", focal_outcome),
        ("Decision Driver", driver or "—"),
    ]
    box = _table(slide, len(box_rows), 2, _MARGIN, _BODY_TOP, 6.6, 2.1)
    box.columns[0].width = Inches(1.9)
    box.columns[1].width = Inches(4.7)
    for r, (label, value) in enumerate(box_rows):
        _set_cell(box.cell(r, 0), label, fill=_LABEL_BG, bold=True, size=10)
        _set_cell(box.cell(r, 1), value, fill=_WHITE, size=10)

    # RAG dimension table
    vendors = ranked
    n = len(vendors)
    dim_w = 2.2
    vendor_w = _clamp((_CONTENT_W - dim_w) / max(n, 1), 0.9, 2.4)
    tbl_w = dim_w + vendor_w * n
    body_sz = 10 if n <= 4 else (9 if n <= 6 else 8)
    table = _table(slide, 1 + len(result.scheme.dimensions), 1 + n, _MARGIN, 3.7, tbl_w,
                   min(0.42 * (len(result.scheme.dimensions) + 1), 3.3))
    table.columns[0].width = Inches(dim_w)
    for j in range(n):
        table.columns[1 + j].width = Inches(vendor_w)
    _header_row(table, ["Dimension"] + [v.vendor for v in vendors], size=body_sz + 1)
    for r, dim in enumerate(result.scheme.dimensions, start=1):
        _set_cell(table.cell(r, 0), dim.name, fill=_LABEL_BG, bold=True, size=body_sz)
        for j, v in enumerate(vendors):
            cell = v.cells.get(dim.id)
            pct = cell.normalized_pct if cell else None
            native = _cell_display(cell)
            label = _adjective(pct) + (f"\n({native})" if native != "—" else "")
            _set_cell(table.cell(r, 1 + j), label, fill=_rag_fill(pct), size=body_sz,
                      align=PP_ALIGN.CENTER)


def _slide_detailed(prs, result: ScorecardResult) -> None:
    slide = _blank_slide(prs)
    _title_block(slide, "Detailed Scoring")
    scheme = result.scheme
    vendors = result.ranked_vendors()
    n = len(vendors)

    show_native = any(v.native_total is not None for v in vendors)
    summary_labels = ["Normalized Total"]
    if show_native:
        tm = scheme.total_max_points
        summary_labels.append("Native Total" + (f" (/{tm:g})" if tm else ""))
    summary_labels.append("Rank")

    n_rows = 1 + len(scheme.dimensions) + len(summary_labels)
    dim_w, notes_w = 2.0, 2.6
    vendor_w = _clamp((_CONTENT_W - dim_w - notes_w) / max(n, 1), 0.7, 2.2)
    tbl_w = dim_w + notes_w + vendor_w * n
    n_cols = 1 + n + 1
    body_sz = 10 if n <= 4 else (9 if n <= 7 else 8)

    table = _table(slide, n_rows, n_cols, _MARGIN, _BODY_TOP, tbl_w, min(0.4 * n_rows, 5.7))
    table.columns[0].width = Inches(dim_w)
    for j in range(n):
        table.columns[1 + j].width = Inches(vendor_w)
    table.columns[n_cols - 1].width = Inches(notes_w)

    _header_row(table, ["Evaluation Dimension"] + [v.vendor for v in vendors] + ["Evaluation Notes"],
                size=body_sz + 1)

    r = 1
    for dim in scheme.dimensions:
        _set_cell(table.cell(r, 0), dim.name, fill=_LABEL_BG, bold=True, size=body_sz)
        for j, v in enumerate(vendors):
            cell = v.cells.get(dim.id)
            prov = cell.provenance if cell else Provenance.NOT_SCORED
            _set_cell(table.cell(r, 1 + j), _cell_display(cell), fill=_FILL.get(prov, _WHITE),
                      size=body_sz, align=PP_ALIGN.CENTER)
        _set_cell(table.cell(r, n_cols - 1), _dim_note(result, dim.id), size=max(body_sz - 1, 7))
        r += 1

    def _summary_value(label: str, v: VendorResult) -> str:
        if label.startswith("Normalized"):
            return f"{v.normalized_total_pct:.0f}%" if v.normalized_total_pct is not None else "—"
        if label.startswith("Native"):
            return f"{v.native_total:.0f}" if v.native_total is not None else "—"
        if label == "Rank":
            return _rank_tag(v)
        return ""

    for label in summary_labels:
        _set_cell(table.cell(r, 0), label, fill=_LABEL_BG, bold=True, size=body_sz)
        for j, v in enumerate(vendors):
            _set_cell(table.cell(r, 1 + j), _summary_value(label, v), fill=_LABEL_BG, bold=True,
                      size=body_sz, align=PP_ALIGN.CENTER)
        note = result.ci.summary if (label == "Rank" and result.ci and result.ci.summary) else ""
        _set_cell(table.cell(r, n_cols - 1), note, fill=_LABEL_BG, size=max(body_sz - 1, 7))
        r += 1


def _dim_note(result: ScorecardResult, dim_id: str) -> str:
    """Short focal-centric CI note for one dimension (detailed slide)."""
    focal = resolve_focal(result)
    scored = [
        (v.vendor, v.cells[dim_id])
        for v in result.vendors
        if dim_id in v.cells and v.cells[dim_id].normalized_pct is not None and not v.disqualified
    ]
    if not scored:
        return ""
    leader_name, leader_cell = max(scored, key=lambda t: t[1].normalized_pct)
    if focal is None or dim_id not in focal.cells or focal.cells[dim_id].normalized_pct is None:
        return f"Leader: {leader_name} ({_cell_display(leader_cell)})."
    fcell = focal.cells[dim_id]
    if fcell.normalized_pct >= leader_cell.normalized_pct - 1e-6:
        base = f"{focal.vendor} leads ({_cell_display(fcell)})."
    else:
        base = f"{focal.vendor} {_cell_display(fcell)} vs leader {leader_name} {_cell_display(leader_cell)}."
    if fcell.rationale:
        base += " " + fcell.rationale.split(". ")[0].strip().rstrip(".") + "."
    return base


def _slide_outcome_drivers(prs, result: ScorecardResult, content: DeckContent) -> None:
    slide = _blank_slide(prs)
    _title_block(slide, "Outcome Drivers")
    drivers = content.drivers

    if drivers is None:
        # graceful fallback: CI summary + weaknesses when no narrative was generated
        summary = result.ci.summary if result.ci else "Competitive analysis unavailable."
        _textbox(slide, summary, _MARGIN, _BODY_TOP, _CONTENT_W, 1.2, size=13)
        if result.ci and result.ci.weaknesses:
            _textbox(slide, "Where " + result.focal_vendor + " trails:", _MARGIN, _BODY_TOP + 1.4,
                     _CONTENT_W, 0.4, size=13, bold=True, color=_IBM_DARK)
            _textbox(slide, _bullets(result.ci.weaknesses[:6]), _MARGIN, _BODY_TOP + 1.85,
                     _CONTENT_W, 2.5, size=12)
        return

    half = (_CONTENT_W - 0.4) / 2

    def _driver_table(title_text, rows, left, header_labels):
        _textbox(slide, title_text, left, _BODY_TOP, half, 0.4, size=15, bold=True, color=_IBM_DARK)
        table = _table(slide, 1 + max(len(rows), 1), 2, left, _BODY_TOP + 0.45, half,
                       min(0.5 * (len(rows) + 1), 4.9))
        table.columns[0].width = Inches(half * 0.42)
        table.columns[1].width = Inches(half * 0.58)
        _header_row(table, header_labels, size=10)
        for r, row in enumerate(rows or [], start=1):
            _set_cell(table.cell(r, 0), row.factor, fill=_LABEL_BG, bold=True, size=9,
                      anchor=MSO_ANCHOR.TOP)
            body = row.evidence + (f"\n→ {row.impact}" if row.impact else "")
            _set_cell(table.cell(r, 1), body, fill=(_ROW_ALT if r % 2 else _WHITE), size=9,
                      anchor=MSO_ANCHOR.TOP)
        if not rows:
            _set_cell(table.cell(1, 0), "—", fill=_LABEL_BG, size=9)
            _set_cell(table.cell(1, 1), "—", fill=_WHITE, size=9)

    _driver_table(f"Why {drivers.winner} Won", drivers.why_won, _MARGIN,
                  ["Winning Factor", "Evidence / Why it mattered"])
    _driver_table(f"Why {drivers.focal} Lost", drivers.why_focal_lost, _MARGIN + half + 0.4,
                  ["Issue Area", "Evidence / Impact"])


def _slide_category_comparison(prs, content: DeckContent) -> None:
    comp = content.comparison
    if comp is None or not comp.rows:
        return
    slide = _blank_slide(prs)
    _title_block(slide, f"{comp.focal} vs. {comp.winner}")
    rows = comp.rows
    table = _table(slide, 1 + len(rows), 3, _MARGIN, _BODY_TOP, _CONTENT_W, min(1.2 * (len(rows) + 1), 5.7))
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches((_CONTENT_W - 2.2) / 2)
    table.columns[2].width = Inches((_CONTENT_W - 2.2) / 2)
    _header_row(table, ["Category", comp.focal, comp.winner], size=12)
    for r, row in enumerate(rows, start=1):
        fill = _ROW_ALT if r % 2 else _WHITE
        _set_cell(table.cell(r, 0), row.category, fill=_LABEL_BG, bold=True, size=11, anchor=MSO_ANCHOR.TOP)
        _set_cell(table.cell(r, 1), _bullets(row.focal_points), fill=fill, size=10, anchor=MSO_ANCHOR.TOP)
        _set_cell(table.cell(r, 2), _bullets(row.winner_points), fill=fill, size=10, anchor=MSO_ANCHOR.TOP)


# ── entry point ──────────────────────────────────────────────────────────── #
def render_pptx(
    result: ScorecardResult,
    path: str,
    template_path: Optional[str] = None,
    deck_content: Optional[DeckContent] = None,
) -> str:
    content = deck_content or DeckContent()
    # ensure derived metadata (winner, vendor list) exists even with no Bedrock content
    if not content.meta.vendors:
        content.meta.vendors = result.vendor_names()
    if not content.meta.winning_vendor:
        w = resolve_winner(result)
        content.meta.winning_vendor = w.vendor if w else ""

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W)
    prs.slide_height = Inches(_SLIDE_H)

    _slide_title(prs, result, content)
    _slide_overview(prs, result, content)
    _slide_final_scoring(prs, result, content)
    _slide_scoring_overview(prs, result, content)
    _slide_detailed(prs, result)
    _slide_outcome_drivers(prs, result, content)
    _slide_category_comparison(prs, content)

    prs.save(path)
    return path
